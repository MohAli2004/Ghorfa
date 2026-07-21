"""Generate Ghorfa AI presentation (5 min slides + live demo notes)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x0F, 0x2C, 0x3F)
TEAL = RGBColor(0x1A, 0x6B, 0x6B)
ACCENT = RGBColor(0xC4, 0x7B, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF5, 0xF3, 0xEF)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
SOFT = RGBColor(0xE8, 0xEE, 0xF0)


def set_run(run, size=18, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)


def add_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    font="Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_bullet_box(slide, left, top, width, height, items, size=16, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = "•  " + item
        set_run(run, size=size, color=color)
    return box


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


def add_multiline(slide, left, top, width, height, text, size=14, color=MUTED, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        set_run(run, size=size, color=color)
    return box


# ----- Slide 1: Title -----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)
add_bar(s, 0, Inches(6.9), prs.slide_width, Inches(0.6), TEAL)
add_text_box(
    s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1),
    "GHORFA", size=54, bold=True, color=WHITE, font="Georgia",
)
add_text_box(
    s, Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.8),
    "AI-Powered Property Discovery", size=28, color=RGBColor(0xB8, 0xD4, 0xD4),
)
add_text_box(
    s, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.6),
    "Semantic Search  ·  Hybrid Recommendations  ·  Live Demo",
    size=18, color=RGBColor(0x9A, 0xB8, 0xC0),
)
add_text_box(
    s, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
    "Capstone Project Presentation  |  15 minutes", size=14, color=WHITE,
)
add_notes(
    s,
    "SPEAKER (~30 sec):\n"
    "Welcome. Today I present Ghorfa — a Laravel property platform with AI for smarter discovery.\n"
    "Structure: 5 min slides on how AI works, then ~10 min live demo of the website.\n"
    "Focus: semantic search and personalized recommendations.",
)

# ----- Slide 2: Agenda -----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT)
add_bar(s, 0, 0, prs.slide_width, Inches(0.12), TEAL)
add_text_box(
    s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
    "Agenda", size=36, bold=True, color=NAVY, font="Georgia",
)

for i, (title, mins, points, bar_color) in enumerate(
    [
        (
            "Part 1 — Slides",
            "~5 min",
            [
                "Problem & solution",
                "Semantic search",
                "Recommendations",
                "How they work together",
            ],
            TEAL,
        ),
        (
            "Part 2 — Live Demo",
            "~10 min",
            [
                "Browse & filters",
                "Arabic / natural-language search",
                "Recommended sort personalization",
                "Like → see ranking change",
            ],
            ACCENT,
        ),
    ]
):
    left = Inches(0.8) + i * Inches(6.1)
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), Inches(5.6), Inches(5.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = SOFT
    add_bar(s, left, Inches(1.4), Inches(5.6), Inches(0.12), bar_color)
    add_text_box(
        s, left + Inches(0.4), Inches(1.8), Inches(4.8), Inches(0.5),
        title, size=22, bold=True, color=NAVY,
    )
    add_text_box(
        s, left + Inches(0.4), Inches(2.35), Inches(4.8), Inches(0.4),
        mins, size=16, color=bar_color,
    )
    add_bullet_box(
        s, left + Inches(0.35), Inches(3.0), Inches(4.9), Inches(3), points, size=18
    )

add_notes(
    s,
    "SPEAKER (~20 sec):\n"
    "First I explain the AI features briefly. Then I switch to the live site for the main demo.",
)

# ----- Slide 3: Problem -----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT)
add_bar(s, 0, 0, prs.slide_width, Inches(0.12), TEAL)
add_text_box(
    s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
    "The Problem", size=36, bold=True, color=NAVY, font="Georgia",
)
add_text_box(
    s, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5),
    "Traditional property search is limited for Lebanese users", size=18, color=MUTED,
)

problems = [
    ("Exact text match fails", "Arabic vs English place names — SQL LIKE misses bilingual queries"),
    ("No personal ranking", "Newest / price sort ignores what the user actually likes"),
    ("Cold start for new users", "Guests and new accounts get no useful personalization"),
]
for i, (t, d) in enumerate(problems):
    top = Inches(2.0) + i * Inches(1.5)
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.3)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = SOFT
    add_bar(s, Inches(0.8), top, Inches(0.12), Inches(1.3), ACCENT)
    add_text_box(
        s, Inches(1.3), top + Inches(0.25), Inches(10.5), Inches(0.4),
        t, size=20, bold=True, color=NAVY,
    )
    add_text_box(
        s, Inches(1.3), top + Inches(0.7), Inches(10.5), Inches(0.4),
        d, size=16, color=MUTED,
    )

add_notes(
    s,
    "SPEAKER (~45 sec):\n"
    "Users search in Arabic or English. Exact keyword search fails on transliteration.\n"
    "Default sorts (price, newest) are not personalized.\n"
    "We needed smarter discovery without confusing the UI with AI labels.",
)

# ----- Slide 4: Solution -----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT)
add_bar(s, 0, 0, prs.slide_width, Inches(0.12), TEAL)
add_text_box(
    s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
    "Our Solution", size=36, bold=True, color=NAVY, font="Georgia",
)
add_text_box(
    s, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.5),
    "Two AI features that run silently inside normal search & Recommended sort",
    size=17, color=MUTED,
)

solutions = [
    (
        "Semantic Search",
        "Match what the user typed",
        "Bilingual expansion + OpenAI embeddings + keyword scoring",
        TEAL,
    ),
    (
        "Hybrid Recommendations",
        "Rank what the user may like",
        "Collaborative + content + location + popularity",
        ACCENT,
    ),
]
for i, (t, q, d, bar_color) in enumerate(solutions):
    left = Inches(0.8) + i * Inches(6.1)
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), Inches(5.7), Inches(4.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = SOFT
    add_bar(s, left, Inches(2.0), Inches(5.7), Inches(0.15), bar_color)
    add_text_box(
        s, left + Inches(0.4), Inches(2.4), Inches(4.9), Inches(0.5),
        t, size=22, bold=True, color=NAVY,
    )
    add_text_box(
        s, left + Inches(0.4), Inches(3.1), Inches(4.9), Inches(0.6),
        q, size=16, bold=True, color=bar_color,
    )
    add_text_box(
        s, left + Inches(0.4), Inches(4.0), Inches(4.9), Inches(1.5),
        d, size=16, color=MUTED,
    )

add_notes(
    s,
    "SPEAKER (~40 sec):\n"
    "Key distinction for the jury:\n"
    "Semantic search = which properties match this text?\n"
    "Recommendations = which properties fit this user?\n"
    "Both are invisible — user just sees search results and Recommended sort.",
)

# ----- Slide 5: Semantic Search -----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT)
add_bar(s, 0, 0, prs.slide_width, Inches(0.12), TEAL)
add_text_box(
    s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.55),
    "Semantic Search — How It Works", size=32, bold=True, color=NAVY, font="Georgia",
)

steps = [
    ("1", "Expand", "Arabic ↔ English\ndictionary + GPT"),
    ("2", "Keyword", "Score city, title,\naddress, amenities"),
    ("3", "Embed", "Query vector vs\nproperty vectors"),
    ("4", "Rank", "60% keyword +\n40% semantic"),
]
for i, (n, t, d) in enumerate(steps):
    left = Inches(0.6) + i * Inches(3.15)
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.3), Inches(2.95), Inches(3.4)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = SOFT
    circle = s.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(1.05), Inches(1.55), Inches(0.7), Inches(0.7)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    add_text_box(
        s, left + Inches(1.05), Inches(1.65), Inches(0.7), Inches(0.55),
        n, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    add_text_box(
        s, left + Inches(0.15), Inches(2.5), Inches(2.65), Inches(0.45),
        t, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
    )
    add_multiline(
        s, left + Inches(0.15), Inches(3.1), Inches(2.65), Inches(1.3),
        d, size=13, color=MUTED, align=PP_ALIGN.CENTER,
    )

add_multiline(
    s, Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.2),
    "Property embeddings are pre-computed at save time. Only the search query is embedded at runtime.\n"
    "Fallback: if OpenAI is unavailable → standard SQL LIKE search.",
    size=15, color=MUTED,
)

add_notes(
    s,
    "SPEAKER (~60 sec):\n"
    "User types an Arabic or English place name. We expand variants, score keywords,\n"
    "embed the query, compare with stored listing vectors using cosine similarity.\n"
    "Final score blends keyword and semantic. Strong keyword hits get a bonus.\n"
    "If API fails, classic LIKE search still works.",
)

# ----- Slide 6: Recommendations -----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT)
add_bar(s, 0, 0, prs.slide_width, Inches(0.12), TEAL)
add_text_box(
    s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.55),
    "Hybrid Recommendations", size=32, bold=True, color=NAVY, font="Georgia",
)
add_text_box(
    s, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.4),
    "final = 0.40×Collaborative + 0.25×Content + 0.25×Location + 0.10×Popularity",
    size=16, bold=True, color=TEAL,
)

comps = [
    ("40%", "Collaborative", "Users like you\nalso liked these"),
    ("25%", "Content", "Type, city, price,\namenities + embeddings"),
    ("25%", "Location", "Haversine distance\nto search area"),
    ("10%", "Popularity", "Views + likes×3\nnormalized"),
]
for i, (pct, t, d) in enumerate(comps):
    left = Inches(0.6) + i * Inches(3.15)
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.7), Inches(2.95), Inches(3.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = SOFT
    add_text_box(
        s, left + Inches(0.15), Inches(2.0), Inches(2.65), Inches(0.55),
        pct, size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
    )
    add_text_box(
        s, left + Inches(0.15), Inches(2.7), Inches(2.65), Inches(0.45),
        t, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
    )
    add_multiline(
        s, left + Inches(0.15), Inches(3.35), Inches(2.65), Inches(1.4),
        d, size=13, color=MUTED, align=PP_ALIGN.CENTER,
    )

add_multiline(
    s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2),
    "Cold start (no history): filter match + location + popularity + recency.\n"
    "On search: Recommended sort reorders the filtered listings grid — no separate AI section.",
    size=15, color=MUTED,
)

add_notes(
    s,
    "SPEAKER (~60 sec):\n"
    "We track views, clicks, likes, contacts with weights.\n"
    "Collaborative finds similar users. Content matches preferences (with optional embedding blend).\n"
    "Location uses GPS distance. Popularity is a light boost.\n"
    "New users get cold-start ranking. In search, Recommended just reorders results.",
)

# ----- Slide 7: Together -----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT)
add_bar(s, 0, 0, prs.slide_width, Inches(0.12), TEAL)
add_text_box(
    s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.55),
    "How They Work Together", size=32, bold=True, color=NAVY, font="Georgia",
)

flow = [
    ("User actions", "Search, view,\nlike, contact", NAVY),
    ("Behavior layer", "Interactions &\nsearch history", TEAL),
    ("Semantic search", "Match the\nquery text", NAVY),
    ("Recommendations", "Personalize\nthe order", TEAL),
]
for i, (t, d, bg) in enumerate(flow):
    left = Inches(0.5) + i * Inches(3.2)
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(2.85), Inches(2.6)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.fill.background()
    add_text_box(
        s, left + Inches(0.15), Inches(1.9), Inches(2.55), Inches(0.5),
        t, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    add_multiline(
        s, left + Inches(0.15), Inches(2.6), Inches(2.55), Inches(1.1),
        d, size=14, color=RGBColor(0xD0, 0xE4, 0xE4), align=PP_ALIGN.CENTER,
    )
    if i < 3:
        add_text_box(
            s, left + Inches(2.7), Inches(2.4), Inches(0.5), Inches(0.5),
            "→", size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
        )

add_multiline(
    s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.0),
    "Stack: Laravel 11 · PHP · OpenAI text-embedding-3-small · gpt-4o-mini (query expansion)\n\n"
    "Design principles: pre-compute embeddings · cache results · graceful fallback · invisible UX",
    size=16, color=MUTED,
)

add_notes(
    s,
    "SPEAKER (~40 sec):\n"
    "Actions feed the behavior layer. Semantic search ranks by meaning of the query.\n"
    "Recommendations personalize order. OpenAI is optional with fallbacks. Ready for live demo.",
)

# ----- Slide 8: Live Demo -----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)
add_bar(s, 0, Inches(6.9), prs.slide_width, Inches(0.6), TEAL)
add_text_box(
    s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8),
    "Live Demo", size=44, bold=True, color=WHITE, font="Georgia",
)
add_text_box(
    s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(0.5),
    "Website walkthrough  ·  ~10 minutes", size=22, color=RGBColor(0xB8, 0xD4, 0xD4),
)

demo_items = [
    "1.  Home & browse listings",
    "2.  Filters + Recommended sort",
    "3.  Semantic search (Arabic / natural language)",
    "4.  Like properties → refresh Recommended order",
    "5.  Map search (optional) + Q&A",
]
box = s.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(11), Inches(3.2))
tf = box.text_frame
tf.word_wrap = True
for i, item in enumerate(demo_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    run = p.add_run()
    run.text = item
    set_run(run, size=18, color=WHITE)

add_notes(
    s,
    """LIVE DEMO SCRIPT (~10 min) — leave this slide up while you demo

1) HOME (1 min)
   - Open homepage, show listings and Recommended For You if visible.
   - Brief: normal property platform + AI under the hood.

2) SEARCH & FILTERS (2 min)
   - Go to /search. Show filters (price, type, amenities).
   - Point to sort dropdown: Recommended is default.
   - Apply a filter (e.g. Rent / city) so jury sees filters still work.

3) SEMANTIC SEARCH (3 min)  *** highlight ***
   - Type Arabic place name OR a natural phrase.
   - Show relevant results that exact SQL might miss.
   - Say: expansion + embeddings + keyword blend; Recommended sort active.

4) PERSONALIZATION (3 min)  *** highlight ***
   - Open 1–2 properties, like them.
   - Return to search with Recommended sort.
   - Refresh / search again — order reflects likes (collaborative + content).
   - Mention cold start for guests if you show logged-out briefly.

5) OPTIONAL MAP + WRAP (1 min)
   - Quick map search if time.
   - Close: AI is silent, degrades gracefully, improves discovery.
   - Invite questions.

BACKUP if OpenAI fails: still show LIKE search + attribute recommendations / popularity.""",
)

# ----- Slide 9: Thank you -----
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)
add_bar(s, 0, Inches(6.9), prs.slide_width, Inches(0.6), TEAL)
add_text_box(
    s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(1),
    "Thank You", size=48, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.CENTER,
)
add_text_box(
    s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.6),
    "Questions & Discussion", size=24, color=RGBColor(0xB8, 0xD4, 0xD4), align=PP_ALIGN.CENTER,
)
add_text_box(
    s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
    "Ghorfa  ·  Semantic Search  ·  Hybrid Recommendations",
    size=16, color=RGBColor(0x8A, 0xA8, 0xB0), align=PP_ALIGN.CENTER,
)
add_notes(
    s,
    "SPEAKER:\n"
    "Thank the jury. Open for questions.\n"
    "Likely Qs: cold start, why hybrid weights, embedding cost/latency, privacy of interactions, fallback.",
)

out = r"c:\Ghorfa-Project\docs\Ghorfa_AI_Presentation.pptx"
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides))
