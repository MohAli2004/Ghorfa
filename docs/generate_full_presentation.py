"""Generate full Ghorfa project presentation (slides + live demo notes)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x0F, 0x2C, 0x3F)
TEAL = RGBColor(0x1A, 0x6B, 0x6B)
ACCENT = RGBColor(0xC4, 0x7B, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF5, 0xF3, 0xEF)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
SOFT = RGBColor(0xE8, 0xEE, 0xF0)


def set_run(run, size=18, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)


def add_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(
    slide, left, top, width, height, text,
    size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=DARK, space=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space)
        run = p.add_run()
        run.text = "•  " + item
        set_run(run, size=size, color=color)
    return box


def add_multiline(slide, left, top, width, height, text, size=14, color=MUTED, align=PP_ALIGN.LEFT, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        set_run(run, size=size, color=color, bold=bold)
    return box


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


def card(slide, left, top, width, height, accent=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = SOFT
    add_bar(slide, left, top, width, Inches(0.1), accent)
    return shape


def header(slide, title, subtitle=None):
    add_bg(slide, LIGHT)
    add_bar(slide, 0, 0, prs.slide_width, Inches(0.12), TEAL)
    add_text_box(
        slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.55),
        title, size=32, bold=True, color=NAVY, font="Georgia",
    )
    if subtitle:
        add_text_box(
            slide, Inches(0.7), Inches(0.95), Inches(12), Inches(0.4),
            subtitle, size=16, color=MUTED,
        )


# ===================== SLIDES =====================

# 1 Title
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)
add_bar(s, 0, Inches(6.9), prs.slide_width, Inches(0.6), TEAL)
add_text_box(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1), "GHORFA", size=56, bold=True, color=WHITE, font="Georgia")
add_text_box(s, Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.7), "Property Rental & Sale Platform", size=28, color=RGBColor(0xB8, 0xD4, 0xD4))
add_text_box(s, Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.5), "Connecting seekers, landlords, and admins — with smarter discovery", size=18, color=RGBColor(0x9A, 0xB8, 0xC0))
add_text_box(s, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4), "Team Capstone Presentation  |  Slides + Live Demo", size=14, color=WHITE)
add_notes(s, "Welcome. Introduce the team and Ghorfa. Mention: short slides first, then live website demo.")

# 2 Agenda
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Agenda", "What we will cover today")
items = [
    ("01", "What is Ghorfa?"),
    ("02", "Roles: User · Landlord · Admin"),
    ("03", "How we built it (stack + Agile)"),
    ("04", "Problems we solve & community impact"),
    ("05", "Key features"),
    ("06", "Semantic search & recommendations"),
    ("07", "Live demo of the website"),
]
for i, (num, label) in enumerate(items):
    col = i % 2
    row = i // 2
    left = Inches(0.7) + col * Inches(6.2)
    top = Inches(1.6) + row * Inches(1.15)
    card(s, left, top, Inches(5.8), Inches(1.0), TEAL if col == 0 else ACCENT)
    add_text_box(s, left + Inches(0.3), top + Inches(0.25), Inches(1), Inches(0.5), num, size=22, bold=True, color=TEAL if col == 0 else ACCENT)
    add_text_box(s, left + Inches(1.3), top + Inches(0.28), Inches(4.2), Inches(0.5), label, size=18, bold=True, color=NAVY)
add_notes(s, "Walk through agenda quickly (~20 sec).")

# 3 Definition
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "What is Ghorfa?", "Definition of the project")
card(s, Inches(0.7), Inches(1.6), Inches(12), Inches(2.2), TEAL)
add_text_box(
    s, Inches(1.1), Inches(1.9), Inches(11.2), Inches(1.6),
    "Ghorfa is a web platform for renting and buying properties. "
    "It helps people find spaces, landlords list and manage properties, "
    "and admins keep the marketplace trusted and organized.",
    size=20, color=DARK,
)
add_bullets(
    s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(2.5),
    [
        "One place to browse, search, like, review, and request rent/buy",
        "Landlords can publish listings after verification",
        "Admins approve landlords, listings, and manage the platform",
        "Smart discovery: semantic search + personalized recommendations",
    ],
    size=17,
)
add_notes(
    s,
    "Ghorfa = room/space in Arabic context — a digital place to find housing.\n"
    "Not only listings: full flow from discovery to transactions with trust (admin approval).",
)

# 4 Roles overview
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Three Main Roles", "Who uses Ghorfa?")
roles = [
    ("User (Client)", TEAL, [
        "Browse & search properties",
        "Like / favorite listings",
        "Write reviews",
        "Request rent or buy",
        "Get recommendations",
        "Apply to become landlord",
    ]),
    ("Landlord", ACCENT, [
        "List spaces for rent/sale",
        "Manage own properties",
        "Handle requests & payments",
        "Upload listing details & photos",
        "Track dashboard activity",
        "Export request documents",
    ]),
    ("Admin", NAVY, [
        "Approve landlord applications",
        "Approve / reject listings",
        "Manage users & rules",
        "Manage amenities",
        "Oversee transactions",
        "Keep the platform trusted",
    ]),
]
for i, (title, color, bullets) in enumerate(roles):
    left = Inches(0.5) + i * Inches(4.2)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(4.0), Inches(5.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = SOFT
    add_bar(s, left, Inches(1.5), Inches(4.0), Inches(0.7), color)
    add_text_box(s, left + Inches(0.2), Inches(1.6), Inches(3.6), Inches(0.5), title, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_bullets(s, left + Inches(0.25), Inches(2.5), Inches(3.5), Inches(4), bullets, size=15, space=10)
add_notes(
    s,
    "User = seeker (client role). Landlord lists after admin approval. Admin is the trust layer.\n"
    "~1 min on this slide.",
)

# 5 Tech stack
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "How We Built It", "Technology stack")
stack = [
    ("Backend", "Laravel (PHP)\nMVC architecture\nMySQL / SQLite"),
    ("Frontend", "HTML · CSS · JavaScript\nBlade templates\nResponsive UI"),
    ("Integrations", "OpenAI (AI search)\nGoogle Maps\nPDF contracts"),
    ("Tools", "Git · GitHub\nGit Bash\nTeam collaboration"),
]
for i, (t, d) in enumerate(stack):
    left = Inches(0.5) + i * Inches(3.2)
    card(s, left, Inches(1.6), Inches(3.0), Inches(4.8), TEAL if i % 2 == 0 else ACCENT)
    add_text_box(s, left + Inches(0.2), Inches(2.0), Inches(2.6), Inches(0.5), t, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_multiline(s, left + Inches(0.2), Inches(2.8), Inches(2.6), Inches(3), d, size=15, color=MUTED, align=PP_ALIGN.CENTER)
add_notes(s, "Laravel for structure & security. HTML/CSS/JS for UI. AI + Maps for discovery. Git/GitHub for teamwork.")

# 6 Agile
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Working as a Team — Agile Scrum", "Same project, shared repository, sprint-based delivery")
agile = [
    ("Sprints", "Work delivered in short cycles with clear goals each sprint"),
    ("Git + GitHub", "Branches, commits, pull requests — collaborate on one codebase"),
    ("Git Bash", "Daily version control: pull, commit, push, resolve conflicts"),
    ("Roles in Scrum", "Plan → develop → review → demo → improve for next sprint"),
]
for i, (t, d) in enumerate(agile):
    top = Inches(1.55) + i * Inches(1.25)
    card(s, Inches(0.7), top, Inches(12), Inches(1.1), TEAL if i % 2 == 0 else ACCENT)
    add_text_box(s, Inches(1.1), top + Inches(0.15), Inches(3), Inches(0.4), t, size=18, bold=True, color=NAVY)
    add_text_box(s, Inches(4.2), top + Inches(0.2), Inches(8), Inches(0.7), d, size=16, color=MUTED)
add_notes(
    s,
    "Emphasize teamwork: one GitHub repo, parallel work, sprints, reviews.\n"
    "Add your real sprint count / team roles if asked.",
)

# 7 Problems
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Problems We Solve", "Why Ghorfa exists")
problems = [
    ("Hard to find the right place", "Scattered listings, weak search, language barriers (Arabic / English)"),
    ("Trust & quality", "Anyone can post anywhere — we add landlord verification & listing approval"),
    ("No personalization", "Users see the same lists — we rank by preference and behavior"),
    ("Disconnected process", "Browse, contact, and deal steps are fragmented — we connect them in one app"),
]
for i, (t, d) in enumerate(problems):
    col = i % 2
    row = i // 2
    left = Inches(0.7) + col * Inches(6.2)
    top = Inches(1.55) + row * Inches(2.5)
    card(s, left, top, Inches(5.9), Inches(2.2), ACCENT)
    add_text_box(s, left + Inches(0.35), top + Inches(0.4), Inches(5.2), Inches(0.5), t, size=18, bold=True, color=NAVY)
    add_text_box(s, left + Inches(0.35), top + Inches(1.0), Inches(5.2), Inches(0.9), d, size=15, color=MUTED)
add_notes(s, "Focus on real pain: discovery, trust, personalization, end-to-end flow.")

# 8 Community
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "How Ghorfa Helps the Community", "Impact beyond the code")
helps = [
    ("Seekers", "Faster discovery of homes and rooms that match language, location, and taste"),
    ("Landlords", "A clear channel to list spaces, reach tenants/buyers, and manage requests"),
    ("Local market", "More transparent, verified listings instead of informal or unreliable posts"),
    ("Digital access", "Bilingual-friendly search helps Arabic and English speakers equally"),
]
for i, (t, d) in enumerate(helps):
    top = Inches(1.55) + i * Inches(1.25)
    card(s, Inches(0.7), top, Inches(12), Inches(1.1), TEAL)
    add_text_box(s, Inches(1.1), top + Inches(0.3), Inches(2.8), Inches(0.5), t, size=18, bold=True, color=TEAL)
    add_text_box(s, Inches(4.0), top + Inches(0.3), Inches(8.3), Inches(0.6), d, size=16, color=DARK)
add_notes(s, "Keep this human: community trust + accessibility + opportunity for landlords and seekers.")

# 9 Features
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Main Features", "What the platform includes")
feats = [
    ("Listings", "Rent & sale properties with photos, amenities, rules"),
    ("Search & filters", "Location, price, type, amenities — plus Recommended sort"),
    ("Map view", "Explore properties on Google Maps"),
    ("Likes & reviews", "Save favorites and share feedback"),
    ("Transactions", "Rent/buy requests, contracts, offline payment flow"),
    ("Notifications", "Stay updated on applications and deals"),
    ("AI discovery", "Semantic search + hybrid recommendations"),
    ("Admin control", "Approvals, users, amenities, rules"),
]
for i, (t, d) in enumerate(feats):
    col = i % 4
    row = i // 4
    left = Inches(0.5) + col * Inches(3.2)
    top = Inches(1.55) + row * Inches(2.6)
    card(s, left, top, Inches(3.05), Inches(2.35), TEAL if (i % 2 == 0) else ACCENT)
    add_text_box(s, left + Inches(0.2), top + Inches(0.45), Inches(2.65), Inches(0.5), t, size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(s, left + Inches(0.2), top + Inches(1.1), Inches(2.65), Inches(1.0), d, size=13, color=MUTED, align=PP_ALIGN.CENTER)
add_notes(s, "Quick overview — live demo will show these. Mention AI as the smart layer on top.")

# 10 Semantic vs Rec
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "Semantic Search & Recommendations", "Two AI features — different jobs, work together")
card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2), TEAL)
add_text_box(s, Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.5), "Semantic Search", size=22, bold=True, color=NAVY)
add_text_box(s, Inches(0.9), Inches(2.4), Inches(5.2), Inches(0.4), "Match what the user typed", size=15, bold=True, color=TEAL)
add_bullets(
    s, Inches(0.85), Inches(3.0), Inches(5.3), Inches(3.4),
    [
        "Arabic ↔ English query expansion",
        "Keyword scoring on listing fields",
        "OpenAI embeddings + cosine similarity",
        "Hybrid rank: keyword + semantic",
        "Fallback to SQL LIKE if AI is down",
    ],
    size=15,
)
card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.2), ACCENT)
add_text_box(s, Inches(7.2), Inches(1.8), Inches(5.2), Inches(0.5), "Recommendations", size=22, bold=True, color=NAVY)
add_text_box(s, Inches(7.2), Inches(2.4), Inches(5.2), Inches(0.4), "Rank what the user may like", size=15, bold=True, color=ACCENT)
add_bullets(
    s, Inches(7.15), Inches(3.0), Inches(5.3), Inches(3.4),
    [
        "Collaborative filtering (similar users)",
        "Content match (type, city, price…)",
        "Location distance buckets",
        "Popularity (views + likes)",
        "Cold-start for new users",
    ],
    size=15,
)
add_notes(
    s,
    "Say the one-liner: Semantic = match the text. Recommendations = match the user.\n"
    "On search, Recommended sort reorders filtered results.",
)

# 11 Relation + formula
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "How They Relate", "Same platform, shared signals, complementary ranking")
flow = [
    ("User actions", "Search, view,\nlike, contact"),
    ("Behavior data", "Interactions &\nsearch history"),
    ("Semantic search", "Relevant to\nthe query"),
    ("Recommendations", "Personalized\norder"),
]
for i, (t, d) in enumerate(flow):
    left = Inches(0.45) + i * Inches(3.2)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.55), Inches(3.0), Inches(2.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY if i % 2 == 0 else TEAL
    shape.line.fill.background()
    add_text_box(s, left + Inches(0.15), Inches(1.85), Inches(2.7), Inches(0.45), t, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_multiline(s, left + Inches(0.15), Inches(2.5), Inches(2.7), Inches(1.1), d, size=14, color=RGBColor(0xD0, 0xE4, 0xE4), align=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(s, left + Inches(2.85), Inches(2.4), Inches(0.4), Inches(0.5), "→", size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

card(s, Inches(0.7), Inches(4.3), Inches(12), Inches(2.4), ACCENT)
add_text_box(s, Inches(1.1), Inches(4.55), Inches(11.2), Inches(0.4), "Recommendation score (with history)", size=16, bold=True, color=NAVY)
add_text_box(
    s, Inches(1.1), Inches(5.1), Inches(11.2), Inches(0.5),
    "final = 0.40×Collaborative + 0.25×Content + 0.25×Location + 0.10×Popularity",
    size=16, bold=True, color=TEAL,
)
add_text_box(
    s, Inches(1.1), Inches(5.75), Inches(11.2), Inches(0.6),
    "Location uses distance buckets (closer = higher). Filters still apply; AI only reorders results.",
    size=14, color=MUTED,
)
add_notes(s, "Mention location buckets briefly if asked. Filters first, then AI ranking.")

# 12 Live demo
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)
add_bar(s, 0, Inches(6.9), prs.slide_width, Inches(0.6), TEAL)
add_text_box(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8), "Live Demo", size=44, bold=True, color=WHITE, font="Georgia")
add_text_box(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.5), "Website walkthrough  ·  ~10 minutes", size=22, color=RGBColor(0xB8, 0xD4, 0xD4))
demo = [
    "1.  Home — browse as a user",
    "2.  Roles — user actions, landlord listing, admin approval (overview)",
    "3.  Search & filters + Recommended sort",
    "4.  Semantic search (Arabic / natural language)",
    "5.  Like properties → personalized order",
    "6.  Q&A",
]
box = s.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11), Inches(3.5))
tf = box.text_frame
tf.word_wrap = True
for i, item in enumerate(demo):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    run = p.add_run()
    run.text = item
    set_run(run, size=18, color=WHITE)
add_notes(
    s,
    """LIVE DEMO (~10 min) — leave this slide up

1) HOME (1 min): Show platform purpose and listings.
2) ROLES (2 min): Quick path — user like/review; landlord property; admin approval screens if possible.
3) SEARCH (2 min): Filters + Recommended default.
4) SEMANTIC (2–3 min): Arabic/English query that shows smart matching.
5) RECOMMENDATIONS (2 min): Like listings, refresh Recommended order.
6) Close and invite questions.

Backup: if OpenAI fails, still demo filters + recommendation ranking + roles.""",
)

# 13 Thank you
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)
add_bar(s, 0, Inches(6.9), prs.slide_width, Inches(0.6), TEAL)
add_text_box(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(1), "Thank You", size=48, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.6), "Questions & Discussion", size=24, color=RGBColor(0xB8, 0xD4, 0xD4), align=PP_ALIGN.CENTER)
add_text_box(
    s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.5),
    "Ghorfa  ·  Users · Landlords · Admins  ·  Smart Discovery",
    size=16, color=RGBColor(0x8A, 0xA8, 0xB0), align=PP_ALIGN.CENTER,
)
add_notes(s, "Thank the jury. Be ready for questions on Agile, roles, AI, and trust/approval flow.")

out = r"c:\Ghorfa-Project\docs\Ghorfa_Full_Project_Presentation.pptx"
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides))
